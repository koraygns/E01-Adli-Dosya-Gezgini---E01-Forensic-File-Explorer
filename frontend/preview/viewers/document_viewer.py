"""
Belge görüntüleyici: PDF, TXT, Office, ODF.
"""
from __future__ import annotations

import io
import re
from typing import Any

from PyQt6.QtCore import Qt, QPoint
from PyQt6.QtGui import QPixmap, QImage, QKeyEvent, QPainter, QColor, QTextDocument
from PyQt6.QtWidgets import (
    QWidget, QVBoxLayout, QHBoxLayout, QLabel, QScrollArea,
    QLineEdit, QPushButton, QFrame, QTextEdit, QTextBrowser, QSizePolicy,
    QTableWidget, QTableWidgetItem, QTabWidget, QHeaderView, QAbstractItemView,
)

from frontend.preview.media_router import ViewerType
from backend.engine.utils.cancel_token import CancellationToken

# Office/ODF uzantıları
EXT_DOC = {".doc", ".docx"}
EXT_XLS = {".xls", ".xlsx"}
EXT_PPT = {".ppt", ".pptx"}
EXT_ODF = {".odt", ".ods", ".odp"}


def _ext(name: str) -> str:
    if not name or "." not in name:
        return ""
    return "." + name.rsplit(".", 1)[-1].lower()


class DocumentViewer(QWidget):
    """PDF, TXT, Office belgeleri."""

    def __init__(
        self,
        session: Any,
        item: dict[str, Any],
        viewer_type: ViewerType,
        cancel_token: CancellationToken,
        cache: Any,
    ):
        super().__init__()
        _DOC_STYLE = """
            DocumentViewer { background: #f4f6fa; }
            QLabel#docFileName { color: #1a1a1a; font-size: 14px; font-weight: 500; padding: 4px 0; }
            QFrame#documentToolbar {
                background: #ffffff; border: 1px solid #dee2e6; padding: 10px 14px;
            }
            QPushButton { background: #e9ecef; color: #212529; border: 1px solid #dee2e6; padding: 8px 14px; font-size: 12px; }
            QPushButton:hover { background: #dee2e6; }
            QPushButton:pressed { background: #ced4da; }
            QLineEdit#docSearch {
                background: #ffffff; color: #212529; border: 1px solid #ced4da;
                padding: 8px 12px; min-width: 200px; font-size: 13px; selection-background-color: #0d6efd; selection-color: #fff;
            }
            QLineEdit#docSearch:focus { border-color: #0d6efd; }
            QScrollArea { background: #ffffff; border: 1px solid #dee2e6; }
            QTextEdit, QTextBrowser { background: #ffffff; color: #212529; border: none; padding: 12px; font-size: 13px; line-height: 1.5; }
            QTableWidget { background: #ffffff; color: #212529; gridline-color: #dee2e6; border: 1px solid #dee2e6; }
            QTableWidget::item { padding: 6px 10px; }
            QTableWidget::item:selected { background: #0d6efd; color: #fff; }
            QHeaderView::section { background: #f1f3f5; color: #212529; padding: 8px 10px; border: 1px solid #dee2e6; font-weight: 600; }
            QTabWidget::pane { border: 1px solid #dee2e6; background: #ffffff; top: -1px; }
            QTabBar::tab { background: #e9ecef; color: #495057; padding: 8px 16px; margin-right: 2px; border: 1px solid #dee2e6; border-bottom: none; }
            QTabBar::tab:selected { background: #ffffff; color: #212529; border-bottom: 1px solid #ffffff; margin-bottom: -1px; }
            QTabBar::tab:hover:!selected { background: #f1f3f5; }
        """
        self.setStyleSheet(_DOC_STYLE)
        self._session = session
        self._item = item
        self._viewer_type = viewer_type
        self._cancel = cancel_token
        self._cache = cache
        self._fit_mode = True
        self._page_index = 0
        self._page_count = 0
        self._doc = None
        self._pdf_search_hits: list[tuple[int, Any]] = []
        self._pdf_search_index = -1
        self._text_widget: QTextEdit | QTextBrowser | None = None
        self._table_tabs: QTabWidget | None = None
        self._table_widgets: list[QTableWidget] = []
        self._table_search_matches: list[tuple[QTableWidget, int, int]] = []
        self._table_search_index = -1

        layout = QVBoxLayout(self)
        layout.setSpacing(10)
        layout.setContentsMargins(12, 12, 12, 12)

        name = item.get("name") or "—"
        self._name_label = QLabel(name)
        self._name_label.setObjectName("docFileName")
        layout.addWidget(self._name_label)

        toolbar = QFrame()
        toolbar.setObjectName("documentToolbar")
        tly = QHBoxLayout(toolbar)
        tly.setSpacing(10)
        self._page_label = QLabel("")
        self._page_label.setStyleSheet("color: #495057; min-width: 90px; font-size: 12px;")
        self._btn_prev_page = QPushButton("◀ Önceki sayfa")
        self._btn_next_page = QPushButton("Sonraki sayfa ▶")
        for b in (self._btn_prev_page, self._btn_next_page):
            b.setCursor(Qt.CursorShape.PointingHandCursor)
        self._btn_prev_page.clicked.connect(self.prev_page)
        self._btn_next_page.clicked.connect(self.next_page)
        tly.addWidget(self._page_label)
        tly.addWidget(self._btn_prev_page)
        tly.addWidget(self._btn_next_page)
        tly.addSpacing(24)
        search_lbl = QLabel("Ara:")
        search_lbl.setStyleSheet("color: #495057; font-size: 12px;")
        tly.addWidget(search_lbl)
        self._search_edit = QLineEdit()
        self._search_edit.setObjectName("docSearch")
        self._search_edit.setPlaceholderText("Metin ara (Ctrl+F)")
        self._search_edit.returnPressed.connect(self._find_next)
        tly.addWidget(self._search_edit)
        self._btn_find_next = QPushButton("Sonraki")
        self._btn_find_prev = QPushButton("Önceki")
        for b in (self._btn_find_next, self._btn_find_prev):
            b.setCursor(Qt.CursorShape.PointingHandCursor)
        self._btn_find_next.clicked.connect(self._find_next)
        self._btn_find_prev.clicked.connect(self._find_prev)
        tly.addWidget(self._btn_find_next)
        tly.addWidget(self._btn_find_prev)
        tly.addStretch()
        layout.addWidget(toolbar)

        self._scroll = QScrollArea(self)
        self._scroll.setWidgetResizable(True)
        self._scroll.setAlignment(Qt.AlignmentFlag.AlignCenter)
        self._label = QLabel("Yükleniyor…")
        self._label.setAlignment(Qt.AlignmentFlag.AlignCenter)
        self._label.setStyleSheet("color: #6c757d; font-size: 13px;")
        self._scroll.setWidget(self._label)
        layout.addWidget(self._scroll, 1)

        # Yükle
        if viewer_type == ViewerType.PDF:
            self._load_pdf()
        elif viewer_type == ViewerType.TXT:
            self._load_txt()
        elif viewer_type == ViewerType.OFFICE:
            self._load_office()
        else:
            self._label.setText("Bilinmeyen belge türü.")

    def _clear_content_widgets(self) -> None:
        """İçerik widget'larını temizle."""
        if self._text_widget:
            self._text_widget.setParent(None)
            self._text_widget = None
        if self._table_tabs:
            self._table_tabs.setParent(None)
            self._table_tabs = None
        self._table_widgets.clear()
        self._scroll.setWidget(self._label)

    def _set_text_content(self, html: bool, content: str) -> None:
        """Metin/HTML göster."""
        self._clear_content_widgets()
        if html:
            w = QTextBrowser(self)
        else:
            w = QTextEdit(self)
        w.setReadOnly(True)
        if html:
            w.setHtml(content)
        else:
            w.setPlainText(content)
        w.setSizePolicy(QSizePolicy.Policy.Expanding, QSizePolicy.Policy.Expanding)
        self._text_widget = w
        self._scroll.setWidget(w)
        self._page_label.setText("")
        self._btn_prev_page.setVisible(False)
        self._btn_next_page.setVisible(False)

    def _set_table_content(self, sheets: list[tuple[str, list[list[Any]]]]) -> None:
        """Excel/ODS tablo görünümü."""
        self._clear_content_widgets()
        self._page_label.setText("")
        self._btn_prev_page.setVisible(False)
        self._btn_next_page.setVisible(False)
        tabs = QTabWidget(self)
        tabs.setDocumentMode(True)
        self._table_widgets = []
        for sheet_name, rows in sheets:
            table = QTableWidget(self)
            table.setEditTriggers(QAbstractItemView.EditTrigger.NoEditTriggers)
            table.setAlternatingRowColors(True)
            table.horizontalHeader().setSectionResizeMode(QHeaderView.ResizeMode.ResizeToContents)
            table.verticalHeader().setSectionResizeMode(QHeaderView.ResizeMode.ResizeToContents)
            if not rows:
                table.setRowCount(1)
                table.setColumnCount(1)
                table.setItem(0, 0, QTableWidgetItem("(Boş)"))
            else:
                max_cols = max(len(r) for r in rows)
                table.setRowCount(len(rows))
                table.setColumnCount(max_cols)
                for i, row in enumerate(rows):
                    for j, cell in enumerate(row):
                        if j >= max_cols:
                            break
                        val = cell if cell is not None else ""
                        table.setItem(i, j, QTableWidgetItem(str(val)))
            table.horizontalHeader().setStretchLastSection(True)
            self._table_widgets.append(table)
            tabs.addTab(table, sheet_name[:30] + ("…" if len(sheet_name) > 30 else ""))
        self._table_tabs = tabs
        self._scroll.setWidget(tabs)

    def _load_pdf(self) -> None:
        inode = self._item.get("inode")
        if inode is None:
            self._label.setText("Dosya açılamadı.")
            return
        try:
            import fitz
        except ImportError:
            self._label.setText("PDF için PyMuPDF gerekli: pip install PyMuPDF")
            return
        data = self._session.read_file_content(inode, offset=0, max_size=50 * 1024 * 1024)
        if not data or self._cancel.is_cancelled():
            self._label.setText("PDF okunamadı.")
            return
        try:
            self._doc = fitz.open(stream=data, filetype="pdf")
            self._page_count = len(self._doc)
            if self._page_count == 0:
                self._label.setText("PDF boş.")
                return
            self._render_page(0, None)
            self._page_index = 0
            self._update_page_label()
            self._btn_prev_page.setVisible(True)
            self._btn_next_page.setVisible(True)
        except Exception as e:
            self._label.setText(f"PDF açılamadı: {e}")

    def _render_page(self, page_index: int, highlight_rect: Any) -> None:
        if not self._doc or page_index < 0 or page_index >= self._page_count:
            return
        try:
            import fitz
        except ImportError:
            return
        page = self._doc[page_index]
        mat = fitz.Matrix(2.0, 2.0)
        pix = page.get_pixmap(matrix=mat, alpha=False)
        buf = pix.tobytes("raw", "RGB")
        qimg = QImage(buf, pix.width, pix.height, QImage.Format.Format_RGB888)
        if highlight_rect is not None:
            scale = 2.0
            r = highlight_rect
            x0, y0 = int(r.x0 * scale), int(r.y0 * scale)
            x1, y1 = int(r.x1 * scale), int(r.y1 * scale)
            painter = QPainter(qimg)
            painter.setCompositionMode(QPainter.CompositionMode.CompositionMode_Multiply)
            painter.fillRect(x0, y0, max(1, x1 - x0), max(1, y1 - y0), QColor(255, 255, 0, 160))
            painter.end()
        pixmap = QPixmap.fromImage(qimg)
        self._label.setPixmap(pixmap)
        self._label.setMinimumSize(100, 100)

    def _update_page_label(self) -> None:
        self._page_label.setText(f"Sayfa {self._page_index + 1} / {self._page_count}")

    def _run_pdf_search(self) -> None:
        q = (self._search_edit.text() or "").strip()
        self._pdf_search_hits = []
        if not q or not self._doc:
            return
        try:
            import fitz
            for i in range(len(self._doc)):
                if self._cancel.is_cancelled():
                    return
                page = self._doc[i]
                for rect in page.search_for(q):
                    self._pdf_search_hits.append((i, rect))
        except Exception:
            pass

    def _run_table_search(self) -> None:
        q = (self._search_edit.text() or "").strip().lower()
        self._table_search_matches = []
        if not q or not self._table_widgets:
            return
        for tw in self._table_widgets:
            for r in range(tw.rowCount()):
                for c in range(tw.columnCount()):
                    it = tw.item(r, c)
                    if it and q in (it.text() or "").lower():
                        self._table_search_matches.append((tw, r, c))

    def _find_next(self) -> None:
        if self._table_widgets:
            self._run_table_search()
            if not self._table_search_matches:
                return
            self._table_search_index = (self._table_search_index + 1) % len(self._table_search_matches)
            tw, row, col = self._table_search_matches[self._table_search_index]
            tw.setCurrentCell(row, col)
            it = tw.item(row, col)
            if it:
                tw.scrollToItem(it)
            if self._table_tabs:
                self._table_tabs.setCurrentWidget(tw)
            return
        if self._text_widget is not None:
            q = self._search_edit.text() or ""
            if not self._text_widget.find(q):
                self._text_widget.moveCursor(self._text_widget.textCursor().Start)
                self._text_widget.find(q)
            return
        if self._doc is not None:
            self._run_pdf_search()
            if not self._pdf_search_hits:
                return
            self._pdf_search_index = (self._pdf_search_index + 1) % len(self._pdf_search_hits)
            page_idx, rect = self._pdf_search_hits[self._pdf_search_index]
            self._page_index = page_idx
            self._render_page(page_idx, rect)
            self._update_page_label()

    def _find_prev(self) -> None:
        if self._table_widgets:
            self._run_table_search()
            if not self._table_search_matches:
                return
            self._table_search_index = (self._table_search_index - 1) % len(self._table_search_matches)
            tw, row, col = self._table_search_matches[self._table_search_index]
            tw.setCurrentCell(row, col)
            it = tw.item(row, col)
            if it:
                tw.scrollToItem(it)
            if self._table_tabs:
                self._table_tabs.setCurrentWidget(tw)
            return
        if self._text_widget is not None:
            q = self._search_edit.text() or ""
            self._text_widget.find(q, QTextDocument.FindFlag.FindBackward)
            return
        if self._doc is not None:
            self._run_pdf_search()
            if not self._pdf_search_hits:
                return
            self._pdf_search_index = (self._pdf_search_index - 1) % len(self._pdf_search_hits)
            page_idx, rect = self._pdf_search_hits[self._pdf_search_index]
            self._page_index = page_idx
            self._render_page(page_idx, rect)
            self._update_page_label()

    def _load_txt(self) -> None:
        inode = self._item.get("inode")
        if inode is None:
            self._label.setText("Dosya açılamadı.")
            return
        data = self._session.read_file_content(inode, offset=0, max_size=10 * 1024 * 1024)
        if not data or self._cancel.is_cancelled():
            self._label.setText("Metin okunamadı.")
            return
        for enc in ("utf-8", "utf-8-sig", "cp1254", "latin-1"):
            try:
                text = data.decode(enc, errors="replace")
                break
            except Exception:
                continue
        else:
            text = data.decode("utf-8", errors="replace")
        self._set_text_content(False, text[:2_000_000] + ("…" if len(text) > 2_000_000 else ""))

    def _load_office(self) -> None:
        inode = self._item.get("inode")
        name = self._item.get("name") or ""
        if inode is None:
            self._label.setText("Dosya açılamadı.")
            return
        data = self._session.read_file_content(inode, offset=0, max_size=50 * 1024 * 1024)
        if not data or self._cancel.is_cancelled():
            self._label.setText("Dosya okunamadı.")
            return
        ext = _ext(name)
        try:
            if ext in EXT_DOC:
                self._load_word(data, ext == ".doc")
            elif ext in EXT_XLS:
                self._load_excel(data, ext == ".xls")
            elif ext in EXT_PPT:
                self._load_powerpoint(data, ext == ".ppt")
            elif ext in EXT_ODF:
                self._load_odf(data, ext)
            else:
                self._label.setText("Desteklenmeyen Office formatı.")
        except Exception as e:
            self._label.setText(f"Açılamadı: {e}")

    def _load_word(self, data: bytes, is_legacy_doc: bool) -> None:
        if is_legacy_doc:
            try:
                import mammoth
                result = mammoth.convert_to_html(io.BytesIO(data))
                html = result.value or "<p>İçerik çıkarılamadı.</p>"
                if result.messages:
                    html += "<!-- " + "; ".join(str(m) for m in result.messages[:3]) + " -->"
                self._set_text_content(True, html)
            except ImportError:
                self._label.setText("Eski .doc için: pip install mammoth")
            except Exception as e:
                self._set_text_content(False, f"Eski .doc okunamadı: {e}\n\nHam metin denemesi yapılıyor…")
                self._try_ole_text(data)
        else:
            try:
                from docx import Document
                doc = Document(io.BytesIO(data))
                parts = []
                for p in doc.paragraphs:
                    parts.append(p.text)
                for table in doc.tables:
                    for row in table.rows:
                        parts.append("\t".join(c.text for c in row.cells))
                self._set_text_content(False, "\n".join(parts) or "(Boş)")
            except ImportError:
                self._label.setText("DOCX için: pip install python-docx")
            except Exception as e:
                self._label.setText(f"DOCX açılamadı: {e}")

    def _try_ole_text(self, data: bytes) -> None:
        """Eski Office'ten metin çıkar."""
        try:
            import olefile
            ole = olefile.OleFileIO(io.BytesIO(data))
            texts = []
            for stream in ole.listdir():
                if stream[0] == "WordDocument" or (len(stream) > 1 and "Contents" in str(stream)):
                    try:
                        d = ole.openstream(stream).read()
                        t = d.decode("utf-8", errors="replace")
                        if re.search(r"[a-zA-ZğüşıöçĞÜŞİÖÇ]{3,}", t):
                            texts.append(t[:50000])
                    except Exception:
                        pass
            ole.close()
            if texts:
                self._set_text_content(False, "\n\n---\n\n".join(texts)[:500000])
        except ImportError:
            pass
        except Exception:
            pass

    def _load_excel(self, data: bytes, is_legacy_xls: bool) -> None:
        try:
            if is_legacy_xls:
                import xlrd
                wb = xlrd.open_workbook(file_contents=data)
                sheets = []
                for s in wb.sheets():
                    rows = [[s.cell_value(row, c) for c in range(s.ncols)] for row in range(s.nrows)]
                    sheets.append((s.name, rows))
                self._set_table_content(sheets) if sheets else self._set_text_content(False, "(Boş)")
            else:
                import openpyxl
                wb = openpyxl.load_workbook(io.BytesIO(data), read_only=True, data_only=True)
                sheets = []
                for name in wb.sheetnames:
                    sh = wb[name]
                    rows = [list(row) for row in sh.iter_rows(values_only=True)]
                    sheets.append((name, rows))
                wb.close()
                self._set_table_content(sheets) if sheets else self._set_text_content(False, "(Boş)")
        except ImportError:
            self._label.setText("Excel için: pip install openpyxl xlrd")
        except Exception as e:
            self._label.setText(f"Excel açılamadı: {e}")

    def _load_powerpoint(self, data: bytes, is_legacy_ppt: bool) -> None:
        if is_legacy_ppt:
            self._try_ole_text(data)
            if self._text_widget is None:
                self._label.setText("Eski .ppt metin çıkarılamadı. .pptx önerilir veya LibreOffice ile dönüştürün.")
            return
        try:
            from pptx import Presentation
            from pptx.util import Inches, Pt
            prs = Presentation(io.BytesIO(data))
            parts = []
            for i, slide in enumerate(prs.slides):
                parts.append(f"--- Slayt {i + 1} ---")
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        parts.append(shape.text)
                    if shape.has_text_frame:
                        for p in shape.text_frame.paragraphs:
                            parts.append(p.text)
            self._set_text_content(False, "\n".join(parts) or "(Boş)")
        except ImportError:
            self._label.setText("PPTX için: pip install python-pptx")
        except Exception as e:
            self._label.setText(f"PPTX açılamadı: {e}")

    def _load_odf(self, data: bytes, ext: str) -> None:
        try:
            from odf.opendocument import load
            from odf import text, teletype, table
            doc = load(io.BytesIO(data))
            parts = []
            if ext == ".odt":
                for el in doc.getElementsByType(text.P):
                    parts.append(teletype.extractText(el))
                for el in doc.getElementsByType(text.H):
                    parts.append(teletype.extractText(el))
            elif ext == ".ods":
                TABLE_NS = ("urn:oasis:names:tc:opendocument:xmlns:table:1.0", "name")
                sheets = []
                for i, tab in enumerate(doc.getElementsByType(table.Table)):
                    tab_name = (getattr(tab, "attributes", {}).get(TABLE_NS) or f"Sayfa {i + 1}")
                    if isinstance(tab_name, bytes):
                        tab_name = tab_name.decode("utf-8", errors="replace")
                    rows = []
                    for row in tab.getElementsByType(table.TableRow):
                        cells = row.getElementsByType(table.TableCell)
                        rows.append([teletype.extractText(c) for c in cells])
                    sheets.append((tab_name, rows))
                if sheets:
                    self._set_table_content(sheets)
                else:
                    self._set_text_content(False, "\n".join(parts) or "(Boş)")
                return
            elif ext == ".odp":
                for el in doc.getElementsByType(text.P):
                    parts.append(teletype.extractText(el))
                for el in doc.getElementsByType(text.H):
                    parts.append(teletype.extractText(el))
            self._set_text_content(False, "\n".join(parts) or "(Boş)")
        except ImportError:
            self._label.setText("ODF için: pip install odfpy")
        except Exception as e:
            self._label.setText(f"ODF açılamadı: {e}")

    def toggle_fit(self) -> None:
        self._fit_mode = not self._fit_mode

    def next_page(self) -> None:
        if self._page_count and self._page_index < self._page_count - 1:
            self._page_index += 1
            self._render_page(self._page_index, None)
            self._update_page_label()

    def prev_page(self) -> None:
        if self._page_index > 0:
            self._page_index -= 1
            self._render_page(self._page_index, None)
            self._update_page_label()

    def is_click_on_media_content(self, global_pos: QPoint) -> bool:
        """Tıklama içerik üzerinde mi."""
        pt = self.mapFromGlobal(global_pos)
        return self._scroll.geometry().contains(pt)

    def handle_key(self, event: QKeyEvent) -> bool:
        if event.key() == Qt.Key.Key_Space:
            self.toggle_fit()
            return True
        if event.key() == Qt.Key.Key_PageDown:
            self.next_page()
            return True
        if event.key() == Qt.Key.Key_PageUp:
            self.prev_page()
            return True
        if event.key() == Qt.Key.Key_F and event.modifiers() & Qt.KeyboardModifier.ControlModifier:
            self._search_edit.setFocus()
            return True
        return False

    def on_preview_close(self) -> None:
        if self._doc is not None:
            try:
                self._doc.close()
            except Exception:
                pass
            self._doc = None
